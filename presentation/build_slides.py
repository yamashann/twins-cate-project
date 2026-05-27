"""Build the 20-minute presentation as a PPTX.

Compact bullet-point slides. Each bullet is short enough to read in one breath
but specific enough to remind the presenter what to talk about.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

DARK = RGBColor(0x1f, 0x2c, 0x44)
GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xc0, 0x39, 0x2b)
SOFT = RGBColor(0xe8, 0xed, 0xf3)


# ---------- helpers ----------------------------------------------------------

def add_title(slide, text, size=30):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = DARK
    return tb


def add_bullets(slide, bullets, left=0.5, top=1.3, width=12.3, height=6.0,
                size=20, line_spacing=1.35):
    """Add a bullet list. `bullets` is a list of strings or (text, level) tuples.

    Level 0 = top bullet (•), level 1 = sub-bullet (–), level 2 = deeper.
    Prefix with '!' to render in accent color/bold (key takeaway).
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    markers = {0: "•  ", 1: "      –  ", 2: "            ·  "}
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        accent = text.startswith("!")
        if accent:
            text = text[1:].strip()
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = markers.get(level, "•  ") + text
        p.font.size = Pt(size - 2 * level)
        if accent:
            p.font.bold = True
            p.font.color.rgb = ACCENT
        else:
            p.font.color.rgb = DARK if level == 0 else GREY
        p.line_spacing = line_spacing
    return tb


def add_table(slide, rows, left=0.5, top=2.0, width=8, row_h=0.4, font=14, header=True):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(n_rows * row_h),
    )
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font)
                    if header and r == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                    else:
                        run.font.color.rgb = DARK
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
    return table_shape


def add_image(slide, fname, left, top, width=None, height=None):
    p = FIG / fname
    if not p.exists():
        return None
    kwargs = {}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(p), Inches(left), Inches(top), **kwargs)


def add_callout(slide, text, left=0.5, top=6.4, width=12.3, height=0.8,
                size=18, color=None, bold=True):
    if color is None:
        color = ACCENT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.fill.solid()
    tb.fill.fore_color.rgb = SOFT
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return tb


# ────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)

tb = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.6))
p = tb.text_frame.paragraphs[0]
p.text = "The Causal Effect of Birthweight"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER
p2 = tb.text_frame.add_paragraph()
p2.text = "on Infant Mortality"
p2.font.size = Pt(40)
p2.font.bold = True
p2.font.color.rgb = DARK
p2.alignment = PP_ALIGN.CENTER

tb2 = s.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.7))
p = tb2.text_frame.paragraphs[0]
p.text = "Twins Dataset — ATE, CATE, and Validation"
p.font.size = Pt(22)
p.font.color.rgb = GREY
p.alignment = PP_ALIGN.CENTER

tb3 = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0))
p = tb3.text_frame.paragraphs[0]
p.text = "Hannah Yamashita"
p.font.size = Pt(20)
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER
p2 = tb3.text_frame.add_paragraph()
p2.text = "Causal Models in Data Science"
p2.font.size = Pt(16)
p2.font.color.rgb = GREY
p2.alignment = PP_ALIGN.CENTER


# ────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Question
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "The Question")

add_bullets(s, [
    "Low birthweight is a strong predictor of infant mortality",
    "But predictors aren't necessarily causes",
    ("Maternal conditions (diabetes, hypertension, smoking) drive both", 1),
    ("→ the observed association could be mostly confounding", 1),
    "",
    "Why it matters for policy:",
    ("If causal → improving fetal growth could reduce mortality", 1),
    ("If confounded → resources better spent on maternal conditions", 1),
    "",
    "!My question: what's the causal effect of being the heavier twin on one-year mortality?",
], top=1.3, size=20)


# ────────────────────────────────────────────────────────────────────────────
# Slide 3 — Twins Dataset
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "The Twins Dataset")

add_bullets(s, [
    "Source: Louizos et al. (2017), preprocessed from CEVAE",
    "11,984 same-sex twin pairs, both infants under 2 kg",
    "Treatment T: heavier twin (1 = heavier)",
    "Outcome Y: one-year mortality (P(Y=1) = 0.18)",
    "Covariates X: 50 maternal / pregnancy characteristics",
    ("Pregnancy: gestational age, prenatal visits, parity", 1),
    ("Maternal demographics: age, education, race", 1),
    ("Risk indicators: diabetes, hypertension, anemia", 1),
    ("Administrative: state, birth month, year", 1),
    "",
    "!Why Twins is special: both potential outcomes observed → design-based ground truth",
], top=1.2, size=18)


# ────────────────────────────────────────────────────────────────────────────
# Slide 4 — Simulated Setup
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Setting Up the Observational Study")

add_bullets(s, [
    "Both potential outcomes observed → ATE trivially identified",
    "To mimic a real observational study, I simulate confounded assignment:",
], top=1.3, size=20)

tb = s.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.3), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = "p_i = sigmoid(z_i^T w + n_i),    T_i ~ Bernoulli(p_i)"
p.font.size = Pt(24)
p.font.name = "Consolas"
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER

add_bullets(s, [
    "Randomly select one twin per pair as 'treated'",
    "Selection probability depends on standardized covariates z_i",
    ("→ treatment now correlates with mortality predictors (confounding)", 1),
    "Observed Y = selected twin's outcome",
    "Other twin's outcome held back as ground truth for §5",
], top=4.0, size=18)

add_callout(s,
    "The trick: ground truth waiting on the other side. Estimators can be checked directly.",
    top=6.5, size=17)


# ────────────────────────────────────────────────────────────────────────────
# Slide 5 — Overlap and Balance
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Overlap and Balance")

add_image(s, "01_overlap.png", left=0.3, top=1.1, height=4.3)
add_image(s, "01_smd_top20.png", left=6.7, top=1.1, height=4.3)

add_bullets(s, [
    "Propensity overlap (left): both groups span 0–1, concentrate near 0.5 → no trimming",
    "SMDs (right): top imbalances 0.15–0.43 in state, birthplace, race, birth order, education",
], top=5.7, size=15)

add_callout(s,
    "Realistic setup — meaningful confounding to adjust for, but overlap preserved.",
    top=6.55, size=16)


# ────────────────────────────────────────────────────────────────────────────
# Slide 6 — Identification
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Identification Assumptions")

add_image(s, "02_dag.png", left=0.3, top=1.3, height=4.0)

add_bullets(s, [
    "Conditional ignorability:  Y(0), Y(1) ⊥ T | X",
    ("Holds by construction in the simulation", 1),
    "",
    "Positivity:  0 < P(T=1 | X) < 1",
    ("Verified by overlap plot", 1),
    "",
    "Consistency:  Y = T·Y(1) + (1−T)·Y(0)",
    "",
    "SUTVA (pair-level):  within-pair contrast, not single-infant intervention",
], left=6.5, top=1.3, size=16, width=6.5)

add_callout(s,
    "Under (1)–(4):   τ = E[ μ_1(X) − μ_0(X) ]",
    top=6.5, size=20)


# ────────────────────────────────────────────────────────────────────────────
# Slide 7 — ATE Estimators
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "ATE Estimators")

add_bullets(s, [
    "Outcome regression (single-model):",
    ("Fit μ_t(x), average  μ̂_1(X) − μ̂_0(X)", 1),
    "",
    "AIPW (doubly robust):",
    ("Influence function combines outcome + propensity model", 1),
    ("Consistent if EITHER model is correctly specified", 1),
    ("SE = sd(ψ) / √n", 1),
    "",
    "Both use 5-fold cross-fitted nuisances",
    ("μ_t = gradient-boosted trees,  e = L2 logistic", 1),
    ("Cross-fitting → each unit's prediction excludes that unit → valid SE", 1),
], top=1.2, size=18)


# ────────────────────────────────────────────────────────────────────────────
# Slide 8 — CATE Estimators
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "CATE Estimators")

add_bullets(s, [
    "Baselines (known failure modes):",
    ("S-learner: one GBM on [X, T] — shrinks heterogeneity to zero", 1),
    ("T-learner: two arm-specific GBMs — sensitive to arm imbalance", 1),
    "",
    "Double/Debiased ML family (the heavy lifters):",
    ("DR-learner (Kennedy 2023): regress AIPW pseudo-outcome on X", 1),
    ("R-learner (Nie & Wager 2021): minimize Σ(ỹ_i − t̃_i τ(X_i))²", 1),
    ("Causal forest (Wager & Athey 2018): honest splits + pointwise CIs", 1),
    "",
    "All orthogonalize → first-stage nuisance errors only contribute second-order bias",
    ("Matters here because mortality is rare (18%) → fragile propensities", 1),
], top=1.2, size=17)


# ────────────────────────────────────────────────────────────────────────────
# Slide 9 — ATE Results
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "ATE Results")

rows = [
    ["Method", "Estimate", "SE", "95% CI"],
    ["Outcome regression", "−0.0242", "0.00057", "(−0.0253, −0.0230)"],
    ["AIPW (doubly robust)", "−0.0248", "0.00647", "(−0.0375, −0.0122)"],
    ["Within-pair benchmark (truth)", "−0.0252", "0.00292", "—"],
]
add_table(s, rows, left=1.5, top=1.3, width=10.3, row_h=0.5, font=17)

add_bullets(s, [
    "Both estimators within 0.001 of the within-pair truth",
    "AIPW CI covers the truth",
    "OR's SE is from fixed-nuisance bootstrap → understates uncertainty; AIPW SE is trustworthy",
    "OR and AIPW agree → ATE is not driven by one estimator's modeling choices",
], top=3.7, size=18)

add_callout(s,
    "Heavier twin reduces one-year mortality by ~2.4–2.5 percentage points.",
    top=6.5, size=18)


# ────────────────────────────────────────────────────────────────────────────
# Slide 10 — CATE Distributions
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "CATE Distributions")

rows = [
    ["Method", "mean(τ̂)", "sd(τ̂)"],
    ["S-learner", "−0.018", "0.012   ← under"],
    ["T-learner", "−0.025", "0.062"],
    ["DR-learner", "−0.025", "0.168   ← over"],
    ["R-learner", "−0.024", "0.059"],
    ["Causal Forest", "−0.024", "0.022   ← under"],
    ["True ITE", "−0.025", "0.320"],
]
add_table(s, rows, left=0.5, top=1.3, width=6.5, row_h=0.45, font=15)

add_bullets(s, [
    "Lower bound on true sd[τ(X)]:",
    ("Bin into 50 bins by CF prediction", 1),
    ("Take between-bin variance of true ITE", 1),
    ("!Gives sd[τ(X)] ≥ 0.079", 1),
    "",
    "Against 0.079 benchmark:",
    ("S (0.012), CF (0.022) under-disperse", 1),
    ("T (0.062), R (0.059) closest", 1),
    ("DR (0.168) over-disperses — mostly noise", 1),
], left=7.5, top=1.3, size=15, width=5.5)

add_bullets(s, [
    "All five agree on the mean. Spreads disagree wildly.",
    "True ITE sd of 0.32 is mechanical (each ITE ∈ {−1, 0, +1}) — compare to 0.079 lower bound, not raw ITE sd.",
], top=5.5, size=15)


# ────────────────────────────────────────────────────────────────────────────
# Slide 11 — CATE Ranking + Calibration
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "CATE Ranking vs. Ground Truth")

rows = [
    ["Method", "PEHE", "Spearman ρ", "Kendall τ"],
    ["S-learner", "0.320", "0.020", "0.016"],
    ["T-learner", "0.318", "0.103", "0.084"],
    ["DR-learner", "0.360", "0.036", "0.029"],
    ["R-learner", "0.321", "0.046", "0.037"],
    ["Causal Forest", "0.316", "0.197", "0.160"],
]
add_table(s, rows, left=0.5, top=1.2, width=6.5, row_h=0.42, font=14)

add_image(s, "06_calibration.png", left=7.5, top=1.15, height=4.2)

add_bullets(s, [
    "PEHE gap is small — binary outcomes are noisy (CF beats baseline by only 1.2%)",
    "Rank correlation is the decision-relevant metric",
    "!CF Spearman ρ = 0.20 → ~2× T-learner, ~5× DR, ~10× S",
    "Calibration: nearly monotone, but slope > 1 at extremes (honest-forest shrinkage)",
], top=5.5, size=15)


# ────────────────────────────────────────────────────────────────────────────
# Slide 12 — Headline Finding
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Headline Finding: Subgroup Failure")

rows = [
    ["Group", "n", "AIPW GATE", "Truth GATE"],
    ["Q1 (shortest gestation)", "6,681", "−0.018", "−0.029  ← strongest"],
    ["Q2 (middle gestation)", "4,058", "−0.028", "−0.018"],
    ["Q3 (longest gestation)", "1,245", "−0.053  ← strongest", "−0.025"],
]
add_table(s, rows, left=0.5, top=1.2, width=12.3, row_h=0.45, font=15)

add_bullets(s, [
    "AIPW says: monotone, effect triples from Q1 to Q3",
    "Truth says: non-monotone, Q1 strongest, opposite ordering",
    "!Analyst with only AIPW would publish 'effect triples with longer gestation' — wrong",
    "",
    "Likely cause: Q3 has only 1,245 pairs; AIPW SE in Q3 is ~3× the others",
    "Not a generic AIPW failure: prenatal-care subgroups agree with truth (~−0.025 across all)",
], top=3.4, size=17)


# ────────────────────────────────────────────────────────────────────────────
# Slide 13 — BLP
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "BLP Confirms the Warning")

add_bullets(s, [
    "Regress AIPW pseudo-outcome AND true ITE on 8 clinical effect modifiers (HC3 SEs)",
], top=1.05, size=16)

rows = [
    ["Feature", "AIPW coef [p]", "Truth coef [p]"],
    ["gestat10", "−0.013  [ 0.025 ]   ← 'significant'", "−0.0001  [ 0.964 ]   ← zero"],
    ["chyper", "+0.112  [ 0.130 ]", "+0.044  [ 0.097 ]"],
    ["preterm", "−0.039  [ 0.369 ]", "+0.031  [ 0.065 ]"],
    ["others (5)", "all p > 0.4", "all p > 0.4"],
]
add_table(s, rows, left=1.0, top=1.7, width=11.3, row_h=0.5, font=15)

add_bullets(s, [
    "AIPW BLP looks defensible on its own — right framework, valid SEs, sub-0.05 p-value",
    "Truth BLP shows the finding is spurious",
    "!Nothing inside AIPW flags the false positive — only the external benchmark does",
], top=4.9, size=17)


# ────────────────────────────────────────────────────────────────────────────
# Slide 14 — ATE Robustness
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "ATE Is Robust")

add_bullets(s, [
    "DoWhy-style refutations:",
], top=1.05, size=18)

rows = [
    ["Test", "Original", "Refuted", "Δ"],
    ["Placebo (permute T)", "−0.0248", "−0.0003", "+0.0245  ✓"],
    ["Random common cause", "−0.0248", "−0.0252", "−0.0003  ✓"],
    ["80% subset", "−0.0248", "−0.0235", "+0.0013  ✓"],
]
add_table(s, rows, left=0.7, top=1.6, width=11.9, row_h=0.45, font=14)

add_bullets(s, [
    "Sensitivity to a simulated unmeasured confounder U:",
    ("Simulate binary U with associations of strength k_T to T and k_Y to Y", 1),
    ("Append U to X, recompute AIPW across a 5×5 grid", 1),
    ("!All 25 cells stay negative (range −0.022 to −0.067)", 1),
    ("No (k_T, k_Y) combination flips the sign of the ATE", 1),
], top=4.0, size=16)


# ────────────────────────────────────────────────────────────────────────────
# Slide 15 — Takeaways
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Takeaways")

add_bullets(s, [
    "!ATE: heavier twin → ~2.4–2.5 pp lower one-year mortality",
    ("Stable across OR, AIPW, truth, refutations, sensitivity grid", 1),
    "",
    "!CATE: causal forest ranks pairs well, but subgroup claims fail the benchmark",
    ("gestat10 is 'significant' in AIPW BLP (p = 0.025), zero in truth (p = 0.96)", 1),
], top=1.1, size=18)

# Lesson callout
tb = s.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.4))
tb.fill.solid()
tb.fill.fore_color.rgb = SOFT
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Methodological lesson:"
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = GREY
p2 = tf.add_paragraph()
p2.text = "AIPW ATE and AIPW subgroup effects can be trustworthy and untrustworthy on the SAME dataset."
p2.font.size = Pt(17)
p2.font.bold = True
p2.font.color.rgb = ACCENT
p3 = tf.add_paragraph()
p3.text = "Nothing inside AIPW flags the failure — only an external benchmark does."
p3.font.size = Pt(17)
p3.font.bold = True
p3.font.color.rgb = ACCENT

add_bullets(s, [
    "Recommended pipeline for similar studies:",
    ("ATE: AIPW + ≥1 auxiliary estimator (e.g. OR) for triangulation", 1),
    ("CATE: causal forest + ≥1 orthogonalized meta-learner for cross-check", 1),
    ("Subgroups only with a design-based oracle OR subgroup propensity diagnostics + TMLE", 1),
], top=5.1, size=16)


# ────────────────────────────────────────────────────────────────────────────
out = Path(__file__).parent / "Twins CATE Presentation.pptx"
prs.save(out)
print(f"Wrote {out}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
